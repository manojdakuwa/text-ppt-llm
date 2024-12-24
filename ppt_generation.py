import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from diffusers import StableDiffusionPipeline
from transformers import pipeline
from PIL import Image
from io import BytesIO
import requests
import base64


# Load the pre-trained DALL-E Mini model and processor 
# model = DalleBartForConditionalGeneration.from_pretrained("flax-community/dalle-mini") 
# processor = DalleBartProcessor.from_pretrained("flax-community/dalle-mini")
# cache_dir = 'pytorch/stable_diffusion_model_cache'
def generate_image(text): 
    try:
        # if torch.cuda.is_available(): 
        #     device = "cuda" 
        #     print("CUDA is available. Using GPU.") 
        # else: 
        #     device = "cpu" 
        #     print("CUDA is not available. Using CPU.")
        # inputs = processor([text], return_tensors="pt") 
        # outputs = model.generate(**inputs, num_return_sequences=1) 
        # generated_images = processor.decode(outputs[0], skip_special_tokens=True) 
        # # Convert the generated image to a PIL image 
        # image = Image.open(BytesIO(requests.get(generated_images[0]).content))
        # Load the model 
        model_id = "CompVis/stable-diffusion-v1-4" 
        pipeline = StableDiffusionPipeline.from_pretrained(model_id)
        image = pipeline(prompt=text, num_inference_steps=50).images[0]
        print("image: ", image)
        image_stream = BytesIO() 
        image.save(image_stream, format='PNG') 
        image_stream.seek(0)
        img_base64 = base64.b64encode(image_stream.getvalue()).decode('utf-8')
        return img_base64
    except Exception as e:
        print("Error generating image", str(e))
        return None

def create_ppt(slides_data, file_name='output_presentation.pptx'):
    # print("slides_data: ", slides_data)
    prs = Presentation()
    try:
        for slide_content in slides_data:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text_frame.paragraphs[0].font.size = Pt(2)
            content = slide.placeholders[1]
            title.text = slide_content['title'] 
            content.text = slide_content['text']
            for paragraph in content.text_frame.paragraphs: 
                paragraph.font.size = Pt(14)
            image = generate_image(slide_content['title'])  # Replace with paid actual image generation function if needed
            # Decode base64 image to BytesIO object 
            img_data = base64.b64decode(image) 
            img_bytes_io = BytesIO(img_data)
            slide.shapes.add_picture(img_bytes_io,Inches(1), Inches(2), width=Inches(2))

        file_name = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_presentation.pptx")
        prs.save(file_name)
        return "200"
    except Exception as e:
        print("Error creating presentation", str(e))
        return "404"

def create_slide_data(response):
    try:
        return [ {'title': f'{item["title"]}', 'text': item["summarize_text"]} for i, item in enumerate(response)]
    except Exception as e:
        print("Error creating slide")

def generate_ppt_csv(grouped_data, charts): 
    prs = Presentation()
    try:
        for context, group in grouped_data.groupby('cluster'): 
            print("context: ", context)
            print("group; ", group)
            slide = prs.slides.add_slide(prs.slide_layouts[5]) 
            title = slide.shapes.title 
            title.text = f"Context {str(context).capitalize()}" # Summary text 
            title.text_frame.paragraphs[0].font.size = Pt(14)
            summary_text = ' '.join(group['product_details'].tolist()) 
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(2)) 
            tf = textbox.text_frame 
            tf.text = summary_text # Add chart image if exists 
            p = tf.add_paragraph() 
            p.text = tf.text 
            p.font.size = Pt(12)
            if context in charts: 
                slide.shapes.add_picture(charts[context], Inches(1), Inches(3.5), Inches(8.5), Inches(3.5)) 
            filename = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_presentation.pptx") 
            prs.save(filename)
        return "200"
    except Exception as e:
        print("Error generating presentation", str(e))
        return "404"

